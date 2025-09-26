from fastapi import APIRouter, HTTPException, UploadFile, File
from fastapi.responses import Response
from typing import Optional, List
import structlog
import re
import os
import tempfile
import numpy as np

# Updated imports to include models and agents for all features
from app.models import (
    NotesParseRequest, NotesParseResponse, SummarizeRequest, SummaryResponse,
    QuestionPaperRequest, QuestionPaperResponse
)
from app.agents.academic_agent import academic_agent
from app.agents.question_paper import question_paper_agent
from app.utils.config import settings
from app.utils.pdf_exporter import pdf_exporter

# Conditional imports for heavy libraries
try:
    import cv2
    OPENCV_AVAILABLE = True
except ImportError:
    OPENCV_AVAILABLE = False

try:
    from PIL import Image, ImageEnhance, ImageFilter, ImageOps
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

logger = structlog.get_logger()
router = APIRouter(prefix="", tags=["Academic Assistant"])


def _extract_text_from_powerpoint(file_path: str) -> str:
    """Extract text from PowerPoint files (.pptx)."""
    try:
        from pptx import Presentation
        
        text_content = ""
        prs = Presentation(file_path)
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                text_content += f"\n--- Slide {i} ---\n" + "\n".join(slide_text)
        
        return text_content.strip()
    except ImportError:
        raise Exception("`python-pptx` is required for PowerPoint support. Please install it.")
    except Exception as e:
        raise Exception(f"Failed to extract text from PowerPoint: {str(e)}")


def _deskew_image(image: np.ndarray) -> np.ndarray:
    """Corrects skew in an image using OpenCV."""
    if not OPENCV_AVAILABLE or not settings.enable_advanced_ocr_preprocessing:
        return image
        
    try:
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        gray = cv2.bitwise_not(gray)
        thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)[1]
        coords = np.column_stack(np.where(thresh > 0))
        angle = cv2.minAreaRect(coords)[-1]
        
        if angle < -45:
            angle = -(90 + angle)
        else:
            angle = -angle
            
        if abs(angle) > 0.5:
            (h, w) = image.shape[:2]
            center = (w // 2, h // 2)
            M = cv2.getRotationMatrix2D(center, angle, 1.0)
            rotated = cv2.warpAffine(image, M, (w, h), 
                                     flags=cv2.INTER_CUBIC, 
                                     borderMode=cv2.BORDER_REPLICATE)
            return rotated
            
        return image
    except Exception as e:
        logger.warning("Deskewing failed, returning original image.", error=str(e))
        return image


def _preprocess_image_for_ocr(image: Image.Image) -> Image.Image:
    """Applies preprocessing steps to improve OCR accuracy."""
    if not PIL_AVAILABLE:
        raise ImportError("`Pillow` is required for image processing.")

    try:
        img_np = np.array(image.convert('RGB'))

        if OPENCV_AVAILABLE and settings.enable_advanced_ocr_preprocessing:
            img_np = _deskew_image(img_np)
            gray = cv2.cvtColor(img_np, cv2.COLOR_BGR2GRAY)
            denoised = cv2.fastNlMeansDenoising(gray, None, h=10, templateWindowSize=7, searchWindowSize=21)
            binary = cv2.adaptiveThreshold(denoised, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
                                           cv2.THRESH_BINARY, 11, 2)
            return Image.fromarray(binary)
        else:
            image = image.convert('L')
            image = ImageOps.autocontrast(image)
            image = image.filter(ImageFilter.MedianFilter(size=3))
            return image

    except Exception as e:
        logger.warning("Image preprocessing failed, using original image.", error=str(e))
        return image.convert('L')


def _is_plausible_text(text: str) -> bool:
    """Check if extracted text is plausible."""
    if not text:
        return False
    
    stripped_text = text.strip()
    
    if len(stripped_text) < 15:
        return False
    
    unique_alpha = len(set(c for c in stripped_text.lower() if c.isalpha()))
    if unique_alpha < 5:
        return False
        
    words = re.findall(r'[a-zA-Z]{2,}', stripped_text)
    if len(words) < 3:
        return False
        
    return True


async def _correct_ocr_text_with_llm(text: str) -> str:
    """Uses an LLM to correct common OCR errors."""
    if not settings.enable_ocr_correction or not _is_plausible_text(text):
        return text

    logger.info("Attempting to correct OCR output with LLM...")
    try:
        correction_request = NotesParseRequest(
            content=text,
            extract_keywords=False,
            extract_concepts=False,
            extract_questions=False
        )
        
        corrected_result = await academic_agent.parse_only(correction_request)
        
        if corrected_result and corrected_result.parsed_content and len(corrected_result.parsed_content) > 0.5 * len(text):
            logger.info("LLM correction successful.")
            return corrected_result.parsed_content
        else:
            logger.warning("LLM correction returned invalid content. Returning original text.")
            return text
    except Exception as e:
        logger.error("LLM-based OCR correction failed.", error=str(e))
        return text


async def _extract_text_with_ocr(pdf_content: bytes) -> str:
    """Core OCR pipeline for PDFs."""
    try:
        from pdf2image import convert_from_bytes
        import pytesseract
    except ImportError as e:
        raise Exception(f"OCR library not found: {e}. Please install `pdf2image` and `pytesseract`.")

    # Configure paths for Windows
    if os.name == 'nt':
        tesseract_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        if os.path.exists(tesseract_path):
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
        
        poppler_path = r"C:\poppler-24.02.0\Library\bin"
        if not os.path.exists(poppler_path):
             poppler_path = None
    else:
        poppler_path = None

    images = convert_from_bytes(
        pdf_content,
        dpi=settings.ocr_dpi,
        thread_count=settings.ocr_parallel_workers,
        poppler_path=poppler_path
    )

    if not images:
        raise Exception("Could not convert PDF to images.")

    full_text = ""
    for i, image in enumerate(images):
        logger.info(f"Processing page {i+1}/{len(images)} with OCR...")
        
        processed_image = _preprocess_image_for_ocr(image)
        config = r'--oem 3 --psm 6'
        page_text = pytesseract.image_to_string(processed_image, config=config, lang='eng')
        
        if _is_plausible_text(page_text):
            full_text += f"\n\n--- Page {i+1} ---\n{page_text}"
        else:
            logger.warning(f"Page {i+1} produced implausible text, skipping.")

    if not full_text:
        raise Exception("OCR could not extract any plausible text from the document.")
        
    corrected_text = await _correct_ocr_text_with_llm(full_text)
    return corrected_text.strip()


async def _extract_text_from_file(file_path: str, filename: str) -> str:
    """Extract text from various file formats."""
    file_ext = filename.lower().split('.')[-1]
    
    if file_ext == 'pdf':
        # Try direct text extraction first
        try:
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = "".join(page.extract_text() for page in reader.pages if page.extract_text())
            if _is_plausible_text(text):
                logger.info("Successfully extracted text directly from PDF.")
                return text.strip()
            else:
                logger.info("Direct text extraction yielded implausible text. Falling back to OCR.")
        except Exception as e:
            logger.warning("Direct PDF text extraction failed, falling back to OCR.", error=str(e))

        # Fallback to OCR
        logger.info("Starting OCR pipeline for PDF.")
        with open(file_path, 'rb') as f:
            pdf_content = f.read()
        return await _extract_text_with_ocr(pdf_content)

    elif file_ext == 'docx':
        try:
            from docx import Document
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)
        except ImportError:
            raise Exception("`python-docx` is required for DOCX support.")
            
    elif file_ext == 'pptx':
        return _extract_text_from_powerpoint(file_path)

    elif file_ext == 'txt':
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
            
    else:
        raise Exception(f"Unsupported file format: {file_ext}")


# PARSING ENDPOINTS

@router.post("/parse", response_model=NotesParseResponse)
async def parse_content(request: NotesParseRequest):
    """
    Parse and summarize academic content using unified Ollama Mistral 7B processing.
    
    This endpoint combines parsing and summarizing functionality:
    - Extracts keywords, concepts, and questions as requested
    - Generates a comprehensive summary of the content
    - Returns structured information with summarized content
    """
    try:
        if not request.content or not request.content.strip():
            raise HTTPException(status_code=400, detail="Content cannot be empty.")
            
        if len(request.content) > settings.max_content_length:
            raise HTTPException(status_code=413, detail=f"Content length exceeds maximum of {settings.max_content_length} chars.")
            
        # Use the unified agent that combines parsing and summarizing
        result = await academic_agent.parse_and_summarize(request)
        logger.info(f"Successfully processed content: {len(request.content)} chars.")
        return result
    except Exception as e:
        logger.error("Content parsing and summarization failed.", error=str(e))
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")


@router.post("/parse-file", response_model=NotesParseResponse)
async def parse_file(
    file: UploadFile = File(...),
    extract_keywords: bool = True,
    extract_concepts: bool = True,
    extract_questions: bool = False
):
    """
    Parse and summarize uploaded academic files with enhanced OCR support.
    
    Supported formats: PDF (text & scanned), DOCX, PPTX, TXT
    
    This endpoint:
    - Extracts text from various file formats
    - Uses advanced OCR for scanned documents
    - Parses content to extract structured information
    - Generates a comprehensive summary using Ollama Mistral 7B
    """
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file provided.")

    file_ext = file.filename.lower().split('.')[-1]
    allowed_extensions = ['pdf', 'docx', 'pptx', 'txt']
    if file_ext not in allowed_extensions:
        raise HTTPException(status_code=400, detail=f"Unsupported file type. Allowed: {', '.join(allowed_extensions)}")

    # Save to temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as temp_file:
        content = await file.read()
        if len(content) > settings.MAX_FILE_SIZE:
            raise HTTPException(status_code=413, detail=f"File too large. Max size: {settings.max_file_size_mb}MB")
        temp_file.write(content)
        temp_path = temp_file.name

    try:
        # Extract text from file
        text_content = await _extract_text_from_file(temp_path, file.filename)

        if not text_content or not _is_plausible_text(text_content):
            raise HTTPException(status_code=422, detail="Could not extract meaningful content from the file.")

        # Truncate if necessary
        if len(text_content) > settings.max_content_length:
            text_content = text_content[:settings.max_content_length]
            logger.warning(f"Content truncated to {settings.max_content_length} characters.")

        # Debug output
        logger.info("=== EXTRACTED TEXT ===")
        print(f"\nExtracted {len(text_content)} characters from {file.filename}")
        print("="*80)

        # Process with unified agent
        parse_request = NotesParseRequest(
            content=text_content,
            extract_keywords=extract_keywords,
            extract_concepts=extract_concepts,
            extract_questions=extract_questions
        )
        result = await academic_agent.parse_and_summarize(parse_request)
        
        logger.info(f"Successfully processed {file.filename}")
        return result

    except Exception as e:
        logger.error("File parsing failed.", filename=file.filename, error=str(e))
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")
    finally:
        if 'temp_path' in locals() and os.path.exists(temp_path):
            os.unlink(temp_path)


# SUMMARIZATION ENDPOINTS

@router.post("/summarize", response_model=SummaryResponse)
async def summarize_content(request: SummarizeRequest):
    """
    Summarize academic content using Ollama Mistral 7B.
    
    This endpoint provides flexible summarization with various options:
    - Multiple summary types: comprehensive, bullet_points, abstract
    - Customizable length and focus areas
    - Optimized for academic content
    """
    try:
        if len(request.content) > settings.max_content_length:
            raise HTTPException(
                status_code=413,
                detail=f"Content too long. Maximum {settings.max_content_length} characters allowed"
            )
        
        result = await academic_agent.summarize_only(request)
        logger.info(f"Successfully summarized text: {len(request.content)} characters")
        return result
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Text summarization failed: {e}")
        raise HTTPException(status_code=500, detail=f"Summarization failed: {str(e)}")


@router.post("/summarize-file", response_model=SummaryResponse)
async def summarize_file(
    file: UploadFile = File(...),
    summary_type: str = "comprehensive",
    max_length: int = 500,
    focus_areas: Optional[List[str]] = None
):
    """
    Summarize content from uploaded academic files using Ollama Mistral 7B.
    
    Supported formats: PDF (text & scanned), DOCX, PPTX, TXT
    Enhanced OCR support for handwritten/scanned academic content
    """
    try:
        if not file.filename:
            raise HTTPException(status_code=400, detail="No file provided")
        
        # Check file size
        file_content = await file.read()
        if len(file_content) > settings.MAX_FILE_SIZE:
            raise HTTPException(
                status_code=413, 
                detail=f"File too large. Maximum size: {settings.MAX_FILE_SIZE / (1024*1024):.1f}MB"
            )
        
        # Validate file extension
        allowed_extensions = ['pdf', 'docx', 'pptx', 'txt', 'md']
        file_ext = file.filename.lower().split('.')[-1]
        if file_ext not in allowed_extensions:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type. Allowed: {', '.join(allowed_extensions)}"
            )
        
        # Save file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as temp_file:
            temp_file.write(file_content)
            temp_path = temp_file.name
        
        try:
            # Extract text from file
            if file_ext in ['txt', 'md']:
                text_content = file_content.decode('utf-8')
            else:
                text_content = await _extract_text_from_file(temp_path, file.filename)
            
            if not text_content or len(text_content.strip()) < 10:
                raise HTTPException(
                    status_code=422,
                    detail="No meaningful text could be extracted from the file"
                )
            
            # Validate content length
            if len(text_content) > settings.max_content_length:
                text_content = text_content[:settings.max_content_length]
                logger.warning(f"Content truncated to {settings.max_content_length} characters")
            
            # Create summarization request
            request = SummarizeRequest(
                content=text_content,
                summary_type=summary_type,
                max_length=max_length,
                focus_areas=focus_areas or []
            )
            
            # Process with agent
            result = await academic_agent.summarize_only(request)
            
            logger.info(f"Successfully summarized {file.filename}: {len(text_content)} chars")
            return result
            
        finally:
            # Clean up temp file
            try:
                os.unlink(temp_path)
            except:
                pass
                
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"File summarization failed: {e}")
        raise HTTPException(status_code=500, detail=f"File processing failed: {str(e)}")


# PDF EXPORT ENDPOINTS

@router.post("/parse/export-pdf")
async def export_parse_results_to_pdf(request: NotesParseRequest):
    """
    Parse content and export results to a formatted PDF report.
    
    Returns a downloadable PDF containing:
    - Summary in bullet points
    - Extracted keywords with importance scores
    - Identified concepts with definitions
    - Generated study questions
    - Processing metadata
    """
    try:
        # First, get the parsing results
        result = await academic_agent.parse_and_summarize(request)
        
        if not result.success:
            raise HTTPException(status_code=500, detail=result.message)
        
        # Generate PDF
        pdf_content = pdf_exporter.export_parse_results(result)
        filename = pdf_exporter.generate_filename("content_analysis", "parse")
        
        # Return as downloadable PDF
        return Response(
            content=pdf_content,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f"attachment; filename={filename}",
                "Content-Length": str(len(pdf_content))
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"PDF export failed: {e}")
        raise HTTPException(status_code=500, detail=f"PDF export failed: {str(e)}")


@router.post("/parse-file/export-pdf")
async def export_parse_file_results_to_pdf(
    file: UploadFile = File(...),
    extract_keywords: bool = True,
    extract_concepts: bool = True,
    extract_questions: bool = False
):
    """
    Parse uploaded file and export results to a formatted PDF report.
    
    Combines file processing with PDF generation for a complete workflow.
    """
    try:
        # First process the file (same logic as parse_file)
        if not file.filename:
            raise HTTPException(status_code=400, detail="No file provided.")

        file_ext = file.filename.lower().split('.')[-1]
        allowed_extensions = ['pdf', 'docx', 'pptx', 'txt']
        if file_ext not in allowed_extensions:
            raise HTTPException(status_code=400, detail=f"Unsupported file type. Allowed: {', '.join(allowed_extensions)}")

        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as temp_file:
            content = await file.read()
            if len(content) > settings.MAX_FILE_SIZE:
                raise HTTPException(status_code=413, detail=f"File too large. Max size: {settings.max_file_size_mb}MB")
            temp_file.write(content)
            temp_path = temp_file.name

        try:
            # Extract text from file
            text_content = await _extract_text_from_file(temp_path, file.filename)

            if not text_content or not _is_plausible_text(text_content):
                raise HTTPException(status_code=422, detail="Could not extract meaningful content from the file.")

            if len(text_content) > settings.max_content_length:
                text_content = text_content[:settings.max_content_length]

            # Process with unified agent
            parse_request = NotesParseRequest(
                content=text_content,
                extract_keywords=extract_keywords,
                extract_concepts=extract_concepts,
                extract_questions=extract_questions
            )
            result = await academic_agent.parse_and_summarize(parse_request)
            
            if not result.success:
                raise HTTPException(status_code=500, detail=result.message)
            
            # Generate PDF
            pdf_content = pdf_exporter.export_parse_results(result, file.filename)
            filename = pdf_exporter.generate_filename(file.filename.split('.')[0], "parse")
            
            return Response(
                content=pdf_content,
                media_type="application/pdf",
                headers={
                    "Content-Disposition": f"attachment; filename={filename}",
                    "Content-Length": str(len(pdf_content))
                }
            )

        finally:
            if 'temp_path' in locals() and os.path.exists(temp_path):
                os.unlink(temp_path)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"File parsing and PDF export failed: {e}")
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")


@router.post("/summarize/export-pdf")
async def export_summary_results_to_pdf(request: SummarizeRequest):
    """
    Summarize content and export results to a formatted PDF report.
    
    Returns a downloadable PDF containing:
    - Summary content
    - Key points extracted from summary
    - Processing statistics and metadata
    """
    try:
        # Get summarization results
        result = await academic_agent.summarize_only(request)
        
        if not result.success:
            raise HTTPException(status_code=500, detail=result.message)
        
        # Generate PDF
        pdf_content = pdf_exporter.export_summary_results(result)
        filename = pdf_exporter.generate_filename("content_summary", "summary")
        
        return Response(
            content=pdf_content,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f"attachment; filename={filename}",
                "Content-Length": str(len(pdf_content))
            }
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Summary PDF export failed: {e}")
        raise HTTPException(status_code=500, detail=f"PDF export failed: {str(e)}")


@router.post("/summarize-file/export-pdf")
async def export_summary_file_results_to_pdf(
    file: UploadFile = File(...),
    summary_type: str = "comprehensive",
    max_length: int = 500,
    focus_areas: Optional[List[str]] = None
):
    """
    Summarize uploaded file and export results to a formatted PDF report.
    
    Combines file processing with PDF generation for a complete workflow.
    """
    try:
        # Process file (same logic as summarize_file)
        if not file.filename:
            raise HTTPException(status_code=400, detail="No file provided")
        
        file_content = await file.read()
        if len(file_content) > settings.MAX_FILE_SIZE:
            raise HTTPException(status_code=413, detail="File too large")
        
        allowed_extensions = ['pdf', 'docx', 'pptx', 'txt', 'md']
        file_ext = file.filename.lower().split('.')[-1]
        if file_ext not in allowed_extensions:
            raise HTTPException(status_code=400, detail="Unsupported file type")
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as temp_file:
            temp_file.write(file_content)
            temp_path = temp_file.name
        
        try:
            # Extract text
            if file_ext in ['txt', 'md']:
                text_content = file_content.decode('utf-8')
            else:
                text_content = await _extract_text_from_file(temp_path, file.filename)
            
            if not text_content or len(text_content.strip()) < 10:
                raise HTTPException(status_code=422, detail="No meaningful text extracted")
            
            if len(text_content) > settings.max_content_length:
                text_content = text_content[:settings.max_content_length]
            
            # Summarize
            request = SummarizeRequest(
                content=text_content,
                summary_type=summary_type,
                max_length=max_length,
                focus_areas=focus_areas or []
            )
            
            result = await academic_agent.summarize_only(request)
            
            if not result.success:
                raise HTTPException(status_code=500, detail=result.message)
            
            # Generate PDF
            pdf_content = pdf_exporter.export_summary_results(result, file.filename)
            filename = pdf_exporter.generate_filename(file.filename.split('.')[0], "summary")
            
            return Response(
                content=pdf_content,
                media_type="application/pdf",
                headers={
                    "Content-Disposition": f"attachment; filename={filename}",
                    "Content-Length": str(len(pdf_content))
                }
            )
            
        finally:
            try:
                os.unlink(temp_path)
            except:
                pass
                
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"File summary PDF export failed: {e}")
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")
    
@router.post("/generate-question-paper", response_model=QuestionPaperResponse, tags=["Question Paper Generation"])
async def generate_question_paper(request: QuestionPaperRequest):
    """
    Generate a question paper. Difficulty and title are automatically set by the test type.
    - **CAT-1**: Easy difficulty
    - **CAT-2**: Tough difficulty
    - **FAT**: Medium difficulty
    """
    try:
        result = await question_paper_agent.generate_question_paper(request)
        if not result.success:
            raise HTTPException(status_code=500, detail=result.message)
        return result
    except Exception as e:
        logger.error("Question paper generation endpoint failed", error=str(e))
        raise HTTPException(status_code=500, detail=f"Failed to generate question paper: {str(e)}")


@router.post("/generate-question-paper/export-pdf", tags=["Question Paper Generation"])
async def export_question_paper_to_pdf(request: QuestionPaperRequest):
    """
    Generate a question paper based on the request and immediately export it as a PDF.
    """
    try:
        generation_result = await question_paper_agent.generate_question_paper(request)
        if not generation_result.success or not generation_result.paper_data:
            raise HTTPException(status_code=500, detail=generation_result.message)
        
        pdf_content = pdf_exporter.export_question_paper(generation_result.paper_data)
        filename = pdf_exporter.generate_filename(
            base_name=generation_result.paper_data.metadata.title,
            file_type="question_paper"
        )
        
        return Response(
            content=pdf_content,
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    except Exception as e:
        logger.error("Question paper PDF export endpoint failed", error=str(e))
        raise HTTPException(status_code=500, detail=f"Failed to export question paper PDF: {str(e)}")

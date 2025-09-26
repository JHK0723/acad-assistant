from fastapi import APIRouter, HTTPException, UploadFile, File
from fastapi.responses import JSONResponse
from typing import Optional
import asyncio
import structlog
import re
import os
import tempfile
from app.models import NotesParseRequest, NotesParseResponse, ErrorResponse
from app.agents.notes_parser import notes_parser_agent
from app.utils.config import settings

logger = structlog.get_logger()
router = APIRouter(prefix="/notes", tags=["Notes Parser"])


def _extract_text_from_powerpoint(file_path: str) -> str:
    """Extract text from PowerPoint files"""
    try:
        from pptx import Presentation
        
        text_content = ""
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"\n--- Slide {slide_num} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text.strip() + "\n"
            
            if slide_text.strip() != f"--- Slide {slide_num} ---":
                text_content += slide_text
        
        return text_content.strip()
    except ImportError:
        raise Exception("python-pptx package is required for PowerPoint support. Please install it.")
    except Exception as e:
        raise Exception(f"Failed to extract text from PowerPoint: {str(e)}")


def _preprocess_image_for_ocr_advanced(image):
    """Advanced image preprocessing for optimal OCR results"""
    try:
        # Import OpenCV properly with error handling
        try:
            import cv2
            import numpy as np
            opencv_available = True
        except ImportError:
            opencv_available = False
            
        from PIL import Image, ImageEnhance, ImageFilter, ImageOps
        
        # Convert to PIL Image if needed
        if not isinstance(image, Image.Image):
            image = Image.fromarray(image)
        
        # Convert to grayscale
        if image.mode != 'L':
            image = image.convert('L')
        
        if opencv_available:
            # Advanced OpenCV preprocessing
            img_array = np.array(image)
            
            # Denoise the image
            img_array = cv2.fastNlMeansDenoising(img_array)
            
            # Apply bilateral filter for edge preservation while reducing noise
            img_array = cv2.bilateralFilter(img_array, 9, 75, 75)
            
            # Improve contrast using CLAHE
            clahe = cv2.createCLAHE(clipLimit=3.0, tileGridSize=(8,8))
            img_array = clahe.apply(img_array)
            
            # Apply adaptive threshold for better binarization
            img_array = cv2.adaptiveThreshold(
                img_array, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
            )
            
            # Morphological operations to clean up
            kernel = np.ones((1,1), np.uint8)
            img_array = cv2.morphologyEx(img_array, cv2.MORPH_CLOSE, kernel)
            img_array = cv2.morphologyEx(img_array, cv2.MORPH_OPEN, kernel)
            
            # Convert back to PIL
            image = Image.fromarray(img_array)
        else:
            # Fallback PIL-only preprocessing
            print("  ℹ️  Using PIL-only preprocessing (OpenCV not available)")
            
            # Enhance contrast
            enhancer = ImageEnhance.Contrast(image)
            image = enhancer.enhance(2.0)
            
            # Enhance sharpness
            enhancer = ImageEnhance.Sharpness(image)
            image = enhancer.enhance(1.5)
            
            # Apply unsharp mask
            image = image.filter(ImageFilter.UnsharpMask(radius=2, percent=150, threshold=3))
            
            # Auto-level the image
            image = ImageOps.autocontrast(image)
        
        return image
        
    except Exception as e:
        print(f"  ⚠️  Image preprocessing failed: {e}, using original")
        return image


def _extract_text_from_image_multi_strategy(image, page_num):
    """Extract text using multiple OCR strategies and return the best result"""
    import pytesseract
    
    strategies = [
        {
            'name': 'Academic Mixed Content',
            'config': r'--oem 3 --psm 6 -c preserve_interword_spaces=1'
        },
        {
            'name': 'Document Auto-detect',
            'config': r'--oem 3 --psm 3 -c preserve_interword_spaces=1'
        },
        {
            'name': 'Single Text Block',
            'config': r'--oem 3 --psm 7'
        },
        {
            'name': 'Legacy Engine',
            'config': r'--oem 0 --psm 6'
        },
        {
            'name': 'LSTM Engine',
            'config': r'--oem 1 --psm 6'
        }
    ]
    
    best_text = ""
    best_score = 0
    
    for strategy in strategies:
        try:
            text = pytesseract.image_to_string(image, config=strategy['config'], lang='eng')
            
            # Score the extracted text
            score = _score_extracted_text(text)
            
            print(f"    📊 {strategy['name']}: {score:.2f} score, {len(text)} chars")
            
            if score > best_score and len(text.strip()) > 10:
                best_text = text
                best_score = score
                
        except Exception as e:
            print(f"    ❌ {strategy['name']} failed: {e}")
    
    return best_text, best_score


def _score_extracted_text(text):
    """Score extracted text quality based on multiple factors"""
    if not text or len(text.strip()) < 5:
        return 0.0
    
    score = 0.0
    text = text.strip()
    
    # Factor 1: Character diversity (0-0.3)
    unique_chars = len(set(text.lower()))
    char_diversity = min(unique_chars / 50.0, 1.0) * 0.3
    score += char_diversity
    
    # Factor 2: Word formation (0-0.3)
    words = text.split()
    if words:
        # Check for reasonable word lengths
        avg_word_length = sum(len(word) for word in words) / len(words)
        word_score = min(avg_word_length / 6.0, 1.0) * 0.3
        score += word_score
    
    # Factor 3: Letter-to-noise ratio (0-0.2)
    letters = sum(1 for c in text if c.isalpha())
    total_chars = len(text)
    if total_chars > 0:
        letter_ratio = letters / total_chars
        score += letter_ratio * 0.2
    
    # Factor 4: Reasonable character distribution (0-0.2)
    # Penalize text with too many special characters or repetitive patterns
    special_chars = sum(1 for c in text if not (c.isalnum() or c.isspace() or c in '.,!?;:()[]{}'))
    if total_chars > 0:
        special_ratio = special_chars / total_chars
        if special_ratio < 0.3:  # Good ratio
            score += 0.2
        elif special_ratio < 0.5:  # Acceptable
            score += 0.1
    
    return min(score, 1.0)


def _setup_tesseract_windows():
    """Setup Tesseract and Poppler paths for Windows"""
    paths = {
        'tesseract': None,
        'poppler': None
    }
    
    # Find Tesseract
    tesseract_candidates = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"
    ]
    
    for path in tesseract_candidates:
        if os.path.exists(path):
            paths['tesseract'] = path
            break
    
    # Find Poppler
    import glob
    poppler_candidates = [
        r"C:\Program Files\poppler\Library\bin",
        r"C:\Program Files\poppler\poppler-*\bin",
        r"C:\Users\*\AppData\Local\Microsoft\WinGet\Packages\oschwartz10612.Poppler_*\poppler-*\Library\bin"
    ]
    
    for pattern in poppler_candidates:
        matches = glob.glob(pattern)
        if matches:
            paths['poppler'] = matches[0]
            break
    
    return paths


async def _extract_text_with_ocr_enhanced(pdf_content: bytes, max_pages: int = 50) -> str:
    """Enhanced OCR extraction optimized for academic handwritten content"""
    try:
        from pdf2image import convert_from_bytes
        import pytesseract
        
        print("\n" + "="*80)
        print("🚀 ENHANCED ACADEMIC OCR SYSTEM")
        print("="*80)
        
        # Setup Windows paths
        if os.name == 'nt':
            paths = _setup_tesseract_windows()
            if paths['tesseract']:
                pytesseract.pytesseract.tesseract_cmd = paths['tesseract']
                print(f"✅ Tesseract found: {paths['tesseract']}")
            else:
                print("⚠️  Tesseract path not found - using system PATH")
        
        # Convert PDF to images with optimal settings for handwritten content
        convert_kwargs = {
            'first_page': 1,
            'last_page': max_pages,
            'dpi': 350,  # Optimal DPI for handwritten content
            'thread_count': 2,
            'fmt': 'PNG'
        }
        
        # Add Poppler path if available
        if os.name == 'nt' and paths.get('poppler'):
            convert_kwargs['poppler_path'] = paths['poppler']
            print(f"✅ Poppler found: {paths['poppler']}")
        
        images = convert_from_bytes(pdf_content, **convert_kwargs)
        
        if not images:
            raise Exception("No images could be extracted from PDF")
        
        print(f"📄 Successfully converted {len(images)} pages to images")
        print("🎯 Processing with handwritten-optimized OCR...")
        
        extracted_text = ""
        successful_pages = 0
        total_score = 0
        
        for i, image in enumerate(images, 1):
            print(f"\n📄 Processing page {i}/{len(images)}...")
            
            # Apply advanced preprocessing
            processed_image = _preprocess_image_for_ocr_advanced(image)
            
            # Try multiple OCR strategies
            page_text, score = _extract_text_from_image_multi_strategy(processed_image, i)
            
            if score > 0.3 and len(page_text.strip()) > 20:
                extracted_text += f"\n--- Page {i} ---\n{page_text.strip()}\n"
                successful_pages += 1
                total_score += score
                print(f"  ✅ Page {i}: {len(page_text)} chars, score: {score:.2f}")
            else:
                print(f"  ❌ Page {i}: Poor quality (score: {score:.2f})")
        
        avg_score = total_score / successful_pages if successful_pages > 0 else 0
        
        print(f"\n📊 FINAL RESULTS:")
        print(f"  📈 Success rate: {successful_pages}/{len(images)} pages ({successful_pages/len(images)*100:.1f}%)")
        print(f"  📊 Average quality score: {avg_score:.2f}")
        print(f"  📝 Total extracted text: {len(extracted_text)} characters")
        
        if successful_pages == 0:
            raise Exception(
                "❌ No readable text extracted from any page.\n"
                "Possible causes:\n"
                "- Handwriting too unclear or faded\n"
                "- Scan quality too low (try scanning at 300+ DPI)\n"
                "- Document in non-English language\n"
                "- Document contains only images/diagrams"
            )
        
        if avg_score < 0.4:
            print("⚠️  WARNING: Low overall quality detected. Results may contain errors.")
        
        print("\n🔤 Sample extracted text (first 1000 chars):")
        print("-" * 40)
        print(extracted_text[:1000])
        print("="*80)
        
        return extracted_text.strip()
        
    except ImportError as e:
        missing = "pdf2image" if "pdf2image" in str(e) else "pytesseract"
        raise Exception(f"OCR package {missing} not installed. Please install it.")
    except Exception as e:
        logger.error(f"Enhanced OCR failed: {e}")
        raise Exception(f"OCR extraction failed: {str(e)}")


async def _extract_text_from_file(file_path: str, filename: str) -> str:
    """Extract text from various file formats"""
    file_ext = filename.lower().split('.')[-1]
    
    if file_ext == 'pdf':
        # Try standard text extraction first
        try:
            import PyPDF2
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                
                for page in pdf_reader.pages[:50]:  # Support up to 50 pages
                    page_text = page.extract_text()
                    if page_text.strip():
                        text += page_text + "\n"
                
                if len(text.strip()) > 100:  # If we got substantial text
                    print("✅ Standard PDF text extraction successful")
                    return text.strip()
                else:
                    print("⚠️  Standard PDF extraction insufficient, trying OCR...")
                    
        except Exception as e:
            print(f"⚠️  Standard PDF extraction failed: {e}, trying OCR...")
        
        # Fallback to OCR
        try:
            import pdfplumber
            
            with open(file_path, 'rb') as file:
                with pdfplumber.open(file) as pdf:
                    text = ""
                    for page in pdf.pages[:50]:  # Support up to 50 pages
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                    
                    if len(text.strip()) > 100:
                        print("✅ pdfplumber extraction successful")
                        return text.strip()
                    else:
                        print("⚠️  pdfplumber extraction insufficient, trying OCR...")
        except Exception as e:
            print(f"⚠️  pdfplumber extraction failed: {e}, trying OCR...")
        
        # OCR fallback
        with open(file_path, 'rb') as file:
            pdf_content = file.read()
        return await _extract_text_with_ocr_enhanced(pdf_content, max_pages=50)
    
    elif file_ext in ['ppt', 'pptx']:
        return _extract_text_from_powerpoint(file_path)
    
    elif file_ext in ['docx']:
        try:
            from docx import Document
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except ImportError:
            raise Exception("python-docx package required for DOCX support")
    
    elif file_ext in ['txt']:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    
    else:
        raise Exception(f"Unsupported file format: {file_ext}")


@router.post("/parse-file", response_model=NotesParseResponse)
async def parse_file(
    file: UploadFile = File(...),
    extract_keywords: bool = True,
    extract_concepts: bool = True,
    extract_questions: bool = False
):
    """
    Parse uploaded academic files (PDF, DOCX, PPTX) and extract structured information.
    
    Supports:
    - PDF: Text-based and scanned/handwritten (with advanced OCR)
    - DOCX: Microsoft Word documents  
    - PPTX: PowerPoint presentations
    - TXT: Plain text files
    """
    try:
        # Validate file
        if not file.filename:
            raise HTTPException(status_code=400, detail="No file provided")
        
        # Check file size (increased limit for academic documents)
        file_content = await file.read()
        if len(file_content) > settings.MAX_FILE_SIZE:
            raise HTTPException(
                status_code=413, 
                detail=f"File too large. Maximum size: {settings.MAX_FILE_SIZE / (1024*1024):.1f}MB"
            )
        
        # Validate file extension
        allowed_extensions = ['pdf', 'docx', 'pptx', 'txt']
        file_ext = file.filename.lower().split('.')[-1]
        if file_ext not in allowed_extensions:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type. Allowed: {', '.join(allowed_extensions)}"
            )
        
        # Save file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as temp_file:
            temp_file.write(file_content)
            temp_path = temp_file.name
        
        try:
            # Extract text from file
            text_content = await _extract_text_from_file(temp_path, file.filename)
            
            if not text_content or len(text_content.strip()) < 10:
                raise HTTPException(
                    status_code=422,
                    detail="No meaningful text could be extracted from the file"
                )
            
            # Validate content length
            if len(text_content) > settings.MAX_CONTENT_LENGTH:
                text_content = text_content[:settings.MAX_CONTENT_LENGTH]
                logger.warning(f"Content truncated to {settings.MAX_CONTENT_LENGTH} characters")
            
            # Create request for agent
            parse_request = NotesParseRequest(
                content=text_content,
                extract_keywords=extract_keywords,
                extract_concepts=extract_concepts,
                extract_questions=extract_questions
            )
            
            # Process with agent
            result = await notes_parser_agent(parse_request)
            
            logger.info(f"Successfully processed {file.filename}: {len(text_content)} chars")
            return result
            
        finally:
            # Clean up temp file
            try:
                os.unlink(temp_path)
            except:
                pass
                
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"File parsing failed: {e}")
        raise HTTPException(status_code=500, detail=f"File processing failed: {str(e)}")


@router.post("/parse", response_model=NotesParseResponse)
async def parse_notes(request: NotesParseRequest):
    """Parse notes content and extract structured information"""
    try:
        result = await notes_parser_agent(request)
        logger.info(f"Successfully processed notes: {len(request.content)} characters")
        return result
    except Exception as e:
        logger.error(f"Notes parsing failed: {e}")
        raise HTTPException(status_code=500, detail=f"Notes processing failed: {str(e)}")

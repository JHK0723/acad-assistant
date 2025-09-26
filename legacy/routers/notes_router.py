from fastapi import APIRouter, HTTPException, UploadFile, File
from fastapi.responses import JSONResponse
from typing import Optional
import asyncio
import structlog
import re
import os
import tempfile
import numpy as np
from app.models import NotesParseRequest, NotesParseResponse
from app.agents.notes_parser import notes_parser_agent
from app.utils.config import settings

# Conditional imports for heavy libraries
try:
    import cv2
    OPENCV_AVAILABLE = True
except ImportError:
    OPENCV_AVAILABLE = False

try:
    from PIL import Image, ImageEnhance, ImageFilter, ImageOps
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

logger = structlog.get_logger()
router = APIRouter(prefix="/notes", tags=["Notes Parser"])


def _extract_text_from_powerpoint(file_path: str) -> str:
    """Extract text from PowerPoint files (.pptx)."""
    try:
        from pptx import Presentation
        
        text_content = ""
        prs = Presentation(file_path)
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                text_content += f"\n--- Slide {i} ---\n" + "\n".join(slide_text)
        
        return text_content.strip()
    except ImportError:
        raise Exception("`python-pptx` is required for PowerPoint support. Please install it.")
    except Exception as e:
        raise Exception(f"Failed to extract text from PowerPoint: {str(e)}")


def _deskew_image(image: np.ndarray) -> np.ndarray:
    """Corrects skew in an image using OpenCV."""
    if not OPENCV_AVAILABLE or not settings.enable_advanced_ocr_preprocessing:
        return image
        
    try:
        # Convert image to grayscale for processing
        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        
        # Invert the image, as findContours and minAreaRect work best on white objects on a black background
        gray = cv2.bitwise_not(gray)
        
        # Threshold the image to create a binary image
        thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)[1]
        
        # Find the coordinates of all non-zero pixels
        coords = np.column_stack(np.where(thresh > 0))
        
        # Compute the minimum area rectangle that encloses the points
        angle = cv2.minAreaRect(coords)[-1]
        
        # The `cv2.minAreaRect` function returns angles in the range [-90, 0).
        # We need to adjust the angle to be upright.
        if angle < -45:
            angle = -(90 + angle)
        else:
            angle = -angle
            
        # Only rotate if the skew is significant enough to matter
        if abs(angle) > 0.5:
            (h, w) = image.shape[:2]
            center = (w // 2, h // 2)
            M = cv2.getRotationMatrix2D(center, angle, 1.0)
            
            # Perform the rotation
            rotated = cv2.warpAffine(image, M, (w, h), 
                                     flags=cv2.INTER_CUBIC, 
                                     borderMode=cv2.BORDER_REPLICATE)
            return rotated
            
        return image
    except Exception as e:
        logger.warning("Deskewing failed, returning original image.", error=str(e))
        return image


def _preprocess_image_for_ocr(image: Image.Image) -> Image.Image:
    """Applies a series of preprocessing steps to an image to improve OCR accuracy."""
    if not PIL_AVAILABLE:
        raise ImportError("`Pillow` is required for image processing.")

    try:
        # Convert to a format OpenCV can use
        img_np = np.array(image.convert('RGB'))

        if OPENCV_AVAILABLE and settings.enable_advanced_ocr_preprocessing:
            # 1. Deskew the image
            img_np = _deskew_image(img_np)
            
            # 2. Convert to grayscale for further processing
            gray = cv2.cvtColor(img_np, cv2.COLOR_BGR2GRAY)

            # 3. Denoise
            # fastNlMeansDenoising is effective for Gaussian noise
            denoised = cv2.fastNlMeansDenoising(gray, None, h=10, templateWindowSize=7, searchWindowSize=21)

            # 4. Binarization using adaptive thresholding
            # This is better than a global threshold for images with varying lighting
            binary = cv2.adaptiveThreshold(denoised, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
                                           cv2.THRESH_BINARY, 11, 2)

            return Image.fromarray(binary)
        else:
            # Fallback to simpler PIL-only methods if OpenCV is not available or disabled
            image = image.convert('L') # Grayscale
            image = ImageOps.autocontrast(image)
            image = image.filter(ImageFilter.MedianFilter(size=3))
            return image

    except Exception as e:
        logger.warning("Image preprocessing failed, using original image.", error=str(e))
        # Return original image in a compatible format
        return image.convert('L')


def _is_plausible_text(text: str) -> bool:
    """
    A more lenient check for plausible text, suitable for handwritten notes
    which may be sparse or contain diagrams.
    """
    if not text:
        return False
    
    stripped_text = text.strip()
    
    # Condition 1: Must have a minimum number of characters.
    if len(stripped_text) < 15:
        return False
    
    # Condition 2: Must have a minimum diversity of alphabetic characters.
    # This helps filter out repetitive garbage like "/////////////////"
    unique_alpha = len(set(c for c in stripped_text.lower() if c.isalpha()))
    if unique_alpha < 5:
        return False
        
    # Condition 3: Must contain at least a few word-like structures.
    # A "word" is defined as having at least 2 alphabetic characters.
    words = re.findall(r'[a-zA-Z]{2,}', stripped_text)
    if len(words) < 3:
        return False
        
    return True


async def _correct_ocr_text_with_llm(text: str) -> str:
    """Uses an LLM to correct common OCR errors in a block of text."""
    if not settings.enable_ocr_correction or not _is_plausible_text(text):
        return text

    logger.info("Attempting to correct OCR output with LLM...")
    try:
        # This prompt is crucial for getting the right kind of correction
        correction_request = NotesParseRequest(
            content=text,
            task_prompt=(
                "The following text was extracted from a document using OCR and may contain errors. "
                "Please correct spelling mistakes, merge words that are incorrectly split, fix formatting, "
                "and make the text more readable, while strictly preserving the original meaning and all information. "
                "Do not summarize or add new content. Return ONLY the corrected text."
            ),
            extract_keywords=False,
            extract_concepts=False,
            extract_questions=False,
            summary_length="same" # Hint to the agent not to shorten
        )
        
        corrected_result = await notes_parser_agent.parse_notes(correction_request)
        
        # The corrected text should be in the 'parsed_content' field.
        if corrected_result and corrected_result.parsed_content and len(corrected_result.parsed_content) > 0.5 * len(text):
            logger.info("LLM correction successful.")
            return corrected_result.parsed_content
        else:
            logger.warning("LLM correction returned empty or invalid content. Returning original text.")
            return text
    except Exception as e:
        logger.error("LLM-based OCR correction failed.", error=str(e))
        return text


async def _extract_text_with_ocr(pdf_content: bytes) -> str:
    """
    The core OCR pipeline. Converts PDF to images, preprocesses them,
    extracts text, and optionally corrects it.
    """
    try:
        from pdf2image import convert_from_bytes
        import pytesseract
    except ImportError as e:
        raise Exception(f"OCR library not found: {e}. Please install `pdf2image` and `pytesseract`.")

    # Configure paths for Windows if not in PATH
    if os.name == 'nt':
        tesseract_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        if os.path.exists(tesseract_path):
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
        
        poppler_path = r"C:\poppler-24.02.0\Library\bin"
        if not os.path.exists(poppler_path):
             poppler_path = None # Let pdf2image search in PATH
    else:
        poppler_path = None

    images = convert_from_bytes(
        pdf_content,
        dpi=settings.ocr_dpi,
        thread_count=settings.ocr_parallel_workers,
        poppler_path=poppler_path
    )

    if not images:
        raise Exception("Could not convert PDF to images. The file might be corrupted or empty.")

    full_text = ""
    for i, image in enumerate(images):
        logger.info(f"Processing page {i+1}/{len(images)} with advanced OCR...")
        
        # 1. Preprocess the image for better OCR
        processed_image = _preprocess_image_for_ocr(image)
        
        # 2. Extract text using Tesseract with an optimized configuration
        # PSM 6: Assume a single uniform block of text. Often better for handwritten notes.
        config = r'--oem 3 --psm 6'
        page_text = pytesseract.image_to_string(processed_image, config=config, lang='eng')
        
        # 3. Validate extracted text before appending
        if _is_plausible_text(page_text):
            full_text += f"\n\n--- Page {i+1} ---\n{page_text}"
        else:
            logger.warning(f"Page {i+1} produced implausible text, skipping.")

    if not full_text:
        raise Exception("OCR could not extract any plausible text from the document after processing all pages.")
        
    # 4. (Optional) Use LLM to correct the entire block of extracted text
    corrected_text = await _correct_ocr_text_with_llm(full_text)
    
    return corrected_text.strip()


async def _extract_text_from_file(file_path: str, filename: str) -> str:
    """
    Routes file to the correct text extraction function based on extension.
    For PDFs, it attempts direct text extraction first, then falls back to OCR.
    """
    file_ext = filename.lower().split('.')[-1]
    
    if file_ext == 'pdf':
        # 1. Try fast, direct text extraction for text-based PDFs
        try:
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = "".join(page.extract_text() for page in reader.pages if page.extract_text())
            if _is_plausible_text(text):
                logger.info("Successfully extracted text directly from PDF.")
                return text.strip()
            else:
                logger.info("Direct text extraction yielded implausible text. Falling back to OCR.")
        except Exception as e:
            logger.warning("Direct PDF text extraction failed, falling back to OCR.", error=str(e))

        # 2. Fallback to the full OCR pipeline for scanned/image-based PDFs
        logger.info("Starting advanced OCR pipeline for PDF.")
        with open(file_path, 'rb') as f:
            pdf_content = f.read()
        return await _extract_text_with_ocr(pdf_content)

    elif file_ext == 'docx':
        try:
            from docx import Document
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)
        except ImportError:
            raise Exception("`python-docx` is required for DOCX support.")
            
    elif file_ext == 'pptx':
        return _extract_text_from_powerpoint(file_path)

    elif file_ext == 'txt':
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
            
    else:
        raise Exception(f"Unsupported file format: {file_ext}")


@router.post("/parse-file", response_model=NotesParseResponse)
async def parse_file(
    file: UploadFile = File(...),
    extract_keywords: bool = True,
    extract_concepts: bool = True,
    extract_questions: bool = False
):
    """
    Parses uploaded academic files (PDF, DOCX, PPTX, TXT) and extracts structured information.
    This endpoint uses a robust pipeline including advanced OCR for scanned documents.
    """
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file provided.")

    file_ext = file.filename.lower().split('.')[-1]
    allowed_extensions = ['pdf', 'docx', 'pptx', 'txt']
    if file_ext not in allowed_extensions:
        raise HTTPException(status_code=400, detail=f"Unsupported file type. Allowed: {', '.join(allowed_extensions)}")

    # Save to a temporary file to handle file operations consistently
    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as temp_file:
        content = await file.read()
        if len(content) > settings.MAX_FILE_SIZE:
            raise HTTPException(status_code=413, detail=f"File too large. Max size: {settings.max_file_size_mb}MB")
        temp_file.write(content)
        temp_path = temp_file.name

    try:
        # Main extraction logic
        text_content = await _extract_text_from_file(temp_path, file.filename)

        if not text_content or not _is_plausible_text(text_content):
            raise HTTPException(status_code=422, detail="Could not extract any meaningful content from the file. The document may be empty, image-based without readable text, or corrupted.")

        # Truncate if necessary AFTER extraction and correction
        if len(text_content) > settings.max_content_length:
            text_content = text_content[:settings.max_content_length]
            logger.warning(f"Content truncated to {settings.max_content_length} characters to fit model context limit.")

        # Debug: Print the full extracted text to terminal for inspection
        logger.info("=== FULL EXTRACTED TEXT START ===")
        print("\n" + "="*80)
        print("EXTRACTED TEXT FROM FILE:")
        print("="*80)
        print(text_content)
        print("="*80)
        print(f"Total characters: {len(text_content)}")
        print("="*80 + "\n")
        logger.info("=== FULL EXTRACTED TEXT END ===")

        # Send to the AI agent for final analysis
        parse_request = NotesParseRequest(
            content=text_content,
            extract_keywords=extract_keywords,
            extract_concepts=extract_concepts,
            extract_questions=extract_questions
        )
        result = await notes_parser_agent.parse_notes(parse_request)
        
        logger.info(f"Successfully processed {file.filename}, extracted and analyzed content.")
        return result

    except Exception as e:
        logger.error("File parsing pipeline failed.", filename=file.filename, error=str(e), exc_info=True)
        raise HTTPException(status_code=500, detail=f"An unexpected error occurred during processing: {str(e)}")
    finally:
        # Ensure the temporary file is always cleaned up
        if 'temp_path' in locals() and os.path.exists(temp_path):
            os.unlink(temp_path)


@router.post("/parse", response_model=NotesParseResponse)
async def parse_notes(request: NotesParseRequest):
    """Directly parses a string of text to extract structured information."""
    try:
        if not request.content or not request.content.strip():
            raise HTTPException(status_code=400, detail="Content cannot be empty.")
            
        if len(request.content) > settings.max_content_length:
            raise HTTPException(status_code=413, detail=f"Content length exceeds maximum of {settings.max_content_length} chars.")
            
        result = await notes_parser_agent.parse_notes(request)
        logger.info(f"Successfully processed direct text input: {len(request.content)} chars.")
        return result
    except Exception as e:
        logger.error("Direct text parsing failed.", error=str(e))
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")

from fastapi import APIRouter, HTTPException, UploadFile, File
from fastapi.responses import JSONResponse
from typing import Optional, List
import asyncio
import structlog
import tempfile
import os
from app.models import SummarizeRequest, SummaryResponse, ErrorResponse
from app.agents.summarizer import summarizer_agent
from app.utils.config import settings

logger = structlog.get_logger()
router = APIRouter(prefix="/summarize", tags=["Summarizer"])


async def _extract_text_from_file(file_path: str, filename: str) -> str:
    """Extract text from various file formats - shared with notes router logic"""
    file_ext = filename.lower().split('.')[-1]
    
    if file_ext == 'pdf':
        # Try standard text extraction first
        try:
            import PyPDF2
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                
                for page in pdf_reader.pages[:50]:  # Support up to 50 pages
                    page_text = page.extract_text()
                    if page_text.strip():
                        text += page_text + "\n"
                
                if len(text.strip()) > 100:  # If we got substantial text
                    print("✅ Standard PDF text extraction successful")
                    return text.strip()
                else:
                    print("⚠️  Standard PDF extraction insufficient, trying OCR...")
                    
        except Exception as e:
            print(f"⚠️  Standard PDF extraction failed: {e}, trying OCR...")
        
        # Fallback to OCR - import the enhanced OCR from notes router
        try:
            from app.routers.notes_router import _extract_text_with_ocr_enhanced
            
            with open(file_path, 'rb') as file:
                pdf_content = file.read()
            return await _extract_text_with_ocr_enhanced(pdf_content, max_pages=50)
        except Exception as e:
            raise Exception(f"Both standard and OCR PDF extraction failed: {str(e)}")
    
    elif file_ext in ['ppt', 'pptx']:
        try:
            from pptx import Presentation
            
            text_content = ""
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"\n--- Slide {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text.strip() + "\n"
                
                if slide_text.strip() != f"--- Slide {slide_num} ---":
                    text_content += slide_text
            
            return text_content.strip()
        except ImportError:
            raise Exception("python-pptx package is required for PowerPoint support")
        except Exception as e:
            raise Exception(f"Failed to extract text from PowerPoint: {str(e)}")
    
    elif file_ext in ['docx']:
        try:
            from docx import Document
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text.strip()
        except ImportError:
            raise Exception("python-docx package required for DOCX support")
        except Exception as e:
            raise Exception(f"Failed to extract text from DOCX: {str(e)}")
    
    elif file_ext in ['txt']:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    
    else:
        raise Exception(f"Unsupported file format: {file_ext}")


@router.post("/file", response_model=SummaryResponse)
async def summarize_file(
    file: UploadFile = File(...),
    summary_type: str = "comprehensive",
    max_length: int = 500,
    focus_areas: Optional[List[str]] = None
):
    """
    Summarize content from uploaded academic files
    
    Supported formats: PDF (text & scanned), DOCX, PPTX, TXT
    Enhanced OCR support for handwritten/scanned academic content
    """
    try:
        # Validate file
        if not file.filename:
            raise HTTPException(status_code=400, detail="No file provided")
        
        # Check file size
        file_content = await file.read()
        if len(file_content) > settings.MAX_FILE_SIZE:
            raise HTTPException(
                status_code=413, 
                detail=f"File too large. Maximum size: {settings.MAX_FILE_SIZE / (1024*1024):.1f}MB"
            )
        
        # Validate file extension
        allowed_extensions = ['pdf', 'docx', 'pptx', 'txt', 'md']
        file_ext = file.filename.lower().split('.')[-1]
        if file_ext not in allowed_extensions:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported file type. Allowed: {', '.join(allowed_extensions)}"
            )
        
        # Save file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as temp_file:
            temp_file.write(file_content)
            temp_path = temp_file.name
        
        try:
            # Extract text from file using enhanced methods
            if file_ext in ['txt', 'md']:
                text_content = file_content.decode('utf-8')
            else:
                text_content = await _extract_text_from_file(temp_path, file.filename)
            
            if not text_content or len(text_content.strip()) < 10:
                raise HTTPException(
                    status_code=422,
                    detail="No meaningful text could be extracted from the file"
                )
            
            # Validate content length
            if len(text_content) > settings.MAX_CONTENT_LENGTH:
                text_content = text_content[:settings.MAX_CONTENT_LENGTH]
                logger.warning(f"Content truncated to {settings.MAX_CONTENT_LENGTH} characters")
            
            # Create summarization request
            request = SummarizeRequest(
                content=text_content,
                summary_type=summary_type,
                max_length=max_length,
                focus_areas=focus_areas or []
            )
            
            # Process with agent
            result = await summarizer_agent(request)
            
            logger.info(f"Successfully summarized {file.filename}: {len(text_content)} chars")
            return result
            
        finally:
            # Clean up temp file
            try:
                os.unlink(temp_path)
            except:
                pass
                
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"File summarization failed: {e}")
        raise HTTPException(status_code=500, detail=f"File processing failed: {str(e)}")


@router.post("/", response_model=SummaryResponse)
async def summarize_text(request: SummarizeRequest):
    """
    Summarize provided text content with customizable options
    """
    try:
        # Validate content length
        if len(request.content) > settings.MAX_CONTENT_LENGTH:
            raise HTTPException(
                status_code=413,
                detail=f"Content too long. Maximum {settings.MAX_CONTENT_LENGTH} characters allowed"
            )
        
        # Process with agent
        result = await summarizer_agent(request)
        
        logger.info(f"Successfully summarized text: {len(request.content)} characters")
        return result
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Text summarization failed: {e}")
        raise HTTPException(status_code=500, detail=f"Summarization failed: {str(e)}")


@router.post("/bullet-points")
async def create_bullet_points(request: SummarizeRequest):
    """Create bullet point summary from content"""
    try:
        # Override summary type for bullet points
        bullet_request = SummarizeRequest(
            content=request.content,
            summary_type="bullet_points",
            max_length=request.max_length,
            focus_areas=request.focus_areas
        )
        
        result = await summarizer_agent(bullet_request)
        return result
        
    except Exception as e:
        logger.error(f"Bullet points creation failed: {e}")
        raise HTTPException(status_code=500, detail=f"Bullet points creation failed: {str(e)}")


@router.post("/abstract")
async def create_abstract(request: SummarizeRequest):
    """Create academic abstract from content"""
    try:
        # Override summary type for abstract
        abstract_request = SummarizeRequest(
            content=request.content,
            summary_type="abstract",
            max_length=min(request.max_length, 300),  # Abstracts are typically shorter
            focus_areas=request.focus_areas
        )
        
        result = await summarizer_agent(abstract_request)
        return result
        
    except Exception as e:
        logger.error(f"Abstract creation failed: {e}")
        raise HTTPException(status_code=500, detail=f"Abstract creation failed: {str(e)}")
